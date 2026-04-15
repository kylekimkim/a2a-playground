"""
File Text Extractor MCP Server — 순수 Python 기반 파일 텍스트 추출 MCP 서버.

지원 형식:
  - PDF        : pypdf
  - Word       : python-docx (.docx)
  - Excel      : openpyxl (.xlsx)
  - PowerPoint : python-pptx (.pptx)
  - 텍스트     : 직접 읽기 (.txt, .md, .csv, .json, .xml, .html 등)

실행:
  python server.py
  → SSE 서버: http://0.0.0.0:8100/sse
"""
import os
import logging
from pathlib import Path

from mcp.server.fastmcp import FastMCP

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

MCP_HOST = os.getenv("MCP_HOST", "0.0.0.0")
MCP_PORT = int(os.getenv("MCP_PORT", "8100"))

mcp = FastMCP("File Text Extractor", host=MCP_HOST, port=MCP_PORT)

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".yaml", ".yml", ".toml", ".ini", ".log", ".rst",
}


TESSERACT_CMD = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
TESSDATA_DIR = r"C:\sjkim\tessdata"


def _ocr_pdf(path: Path) -> str:
    """이미지 기반 PDF를 Tesseract OCR로 텍스트 추출."""
    import io
    import pytesseract
    from PIL import Image
    from pypdf import PdfReader

    pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD
    reader = PdfReader(str(path), strict=False)
    results = []
    for i, page in enumerate(reader.pages):
        page_texts = []
        for img_obj in page.images:
            img = Image.open(io.BytesIO(img_obj.data))
            text = pytesseract.image_to_string(
                img, lang="kor+eng",
                config=rf"--tessdata-dir {TESSDATA_DIR}",
            )
            if text.strip():
                page_texts.append(text.strip())
        if page_texts:
            results.append(f"[Page {i + 1}]\n" + "\n".join(page_texts))
    return "\n\n".join(results)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path), strict=False)
    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n".join(pages).strip()
    if text:
        return text
    # 텍스트 레이어 없음 → 이미지 기반 PDF, OCR로 재시도
    logger.info("텍스트 레이어 없음 — OCR 시도 중")
    return _ocr_pdf(path)


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        lines.append(f"[Sheet: {sheet}]")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(c) if c is not None else "" for c in row)
            if row_text.strip():
                lines.append(row_text)
    wb.close()
    return "\n".join(lines)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[Slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)
    return "\n".join(lines)


def _extract_text_file(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp949", "euc-kr"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


@mcp.tool()
def extract_text_from_file(file_path: str) -> str:
    """
    로컬 파일에서 텍스트를 추출합니다.
    PDF, Word(.docx), Excel(.xlsx), PowerPoint(.pptx), 텍스트 파일 등을 지원합니다.

    Args:
        file_path: 텍스트를 추출할 로컬 파일의 절대 경로

    Returns:
        추출된 텍스트 내용. 추출 실패 시 오류 메시지 반환.
    """
    path = Path(file_path)
    if not path.exists():
        return f"오류: 파일을 찾을 수 없습니다 — {file_path}"
    if not path.is_file():
        return f"오류: 경로가 파일이 아닙니다 — {file_path}"

    suffix = path.suffix.lower()
    logger.info(f"텍스트 추출 시작: {file_path} ({suffix})")

    try:
        if suffix == ".pdf":
            content = _extract_pdf(path)
        elif suffix == ".docx":
            content = _extract_docx(path)
        elif suffix in (".xlsx", ".xls"):
            content = _extract_xlsx(path)
        elif suffix == ".pptx":
            content = _extract_pptx(path)
        elif suffix in TEXT_EXTENSIONS:
            content = _extract_text_file(path)
        else:
            return f"지원하지 않는 파일 형식입니다: {suffix}\n지원 형식: PDF, DOCX, XLSX, PPTX, TXT, MD, CSV, JSON, XML, HTML 등"

        content = content.strip()
        if not content:
            return "텍스트를 추출할 수 없습니다 (빈 문서일 수 있습니다)."

        logger.info(f"추출 완료: {len(content)} 글자")
        return content

    except Exception as e:
        logger.error(f"텍스트 추출 실패: {e}")
        return f"텍스트 추출 중 오류 발생: {e}"


@mcp.tool()
def get_file_metadata(file_path: str) -> dict:
    """
    파일의 기본 메타데이터를 조회합니다 (크기, 확장자, PDF의 경우 페이지 수 등).

    Args:
        file_path: 메타데이터를 조회할 로컬 파일의 절대 경로

    Returns:
        메타데이터 딕셔너리.
    """
    path = Path(file_path)
    if not path.exists():
        return {"error": f"파일을 찾을 수 없습니다 — {file_path}"}

    stat = path.stat()
    result = {
        "name": path.name,
        "extension": path.suffix.lower(),
        "size_bytes": stat.st_size,
        "size_kb": round(stat.st_size / 1024, 2),
    }

    if path.suffix.lower() == ".pdf":
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path), strict=False)
            result["page_count"] = len(reader.pages)
            meta = reader.metadata
            if meta and meta.title:
                result["title"] = meta.title
            if meta and meta.author:
                result["author"] = meta.author
        except Exception as e:
            result["pdf_meta_error"] = str(e)

    return result


if __name__ == "__main__":
    logger.info(f"File Text Extractor MCP 서버 시작: http://{MCP_HOST}:{MCP_PORT}/sse")
    mcp.run(transport="sse")
